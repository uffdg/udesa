#!/usr/bin/env python3
"""design-system/components/pptx_kit.py

Kit compartido de helpers `python-pptx` para construir presentaciones
.pptx de CUALQUIER entregable de la maestría, replicando el lenguaje
visual editorial de la fuente canónica:

    materias/mt10-innovacion-tecnologica/fidelizacion-coto/entregable/
    presentacion/plataforma-fidelizacion-coto.pptx

Ver `design-system/components/pptx-pitch-deck.md` para la documentación
completa de cada pieza (medidas, cuándo usar cada una, por qué). Este
módulo es la implementación de ese documento — si cambiás una medida acá,
actualizá también el .md, y viceversa.

Regla dura: cualquier `build-pptx.py` nuevo en este repo importa este
módulo en vez de reinventar textboxes y colores sueltos. Si te falta una
pieza (un patrón nuevo que se repite), agregala acá Y documentala en
pptx-pitch-deck.md — no la dejes enterrada en el build-pptx.py de un solo
entregable.

Convención de unidades (importante, para evitar bugs de shapes fuera de
canvas): TODA coordenada/tamaño en la API pública de este módulo es un
**float en pulgadas** (nunca un objeto `Length`/`Emu`/`Inches(...)` ya
envuelto) — el wrapping a EMU pasa una sola vez, adentro de `textbox()` y
`_shape()`. Si necesitás una constante de layout, importala de acá
(`MARGIN_L`, `CONTENT_W`, etc.) — son floats, no `Length`.

Requiere python-pptx (`pip install python-pptx`).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

# ---------------------------------------------------------------------------
# Paleta — extraída 1:1 de plataforma-fidelizacion-coto.pptx (mt10) con
# python-pptx el 2026-08-04. No aproximar valores a mano; si falta un color
# nuevo, extraerlo del pptx canónico primero.
# ---------------------------------------------------------------------------
BRAND = RGBColor(0x50, 0x2B, 0xD8)        # acento primario: kicker, chip activo, cifras, dashes
BRAND_PALE = RGBColor(0xEC, 0xE7, 0xFA)   # fondo de callout neutro / card destacada
INK = RGBColor(0x22, 0x22, 0x22)          # texto principal; también fondo de título/cierre
MUTED = RGBColor(0x6E, 0x6E, 0x6E)        # texto secundario, footer, chip inactivo, labels
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF8, 0xF8)     # fondo full-bleed de cualquier slide de contenido
BORDER = RGBColor(0xED, 0xED, 0xED)       # líneas divisorias de tabla, filas alternadas
ORANGE = RGBColor(0xFA, 0x5D, 0x29)       # acento de costo / riesgo / urgencia negativa
ORANGE_PALE = RGBColor(0xFD, 0xE6, 0xDA)  # fondo de callout de riesgo / spotlight de riesgo
DARK = RGBColor(0x22, 0x22, 0x22)         # fondo de slide de título / cierre
DARK_PURPLE = RGBColor(0x33, 0x2A, 0x52)  # fondo del callout "qué pedimos" en el cierre
ON_DARK_SUBTLE = RGBColor(0xAF, 0xAF, 0xAF)          # subtítulo sobre fondo oscuro
ON_DARK_FAINT = RGBColor(0x8A, 0x8A, 0x8A)           # línea de créditos, la más tenue
ON_DARK_TINT = RGBColor(0x9A, 0x8F, 0xC9)            # meta-línea sobre fondo oscuro (tinte de marca)
ON_DARK_CALLOUT_KICKER = RGBColor(0xC4, 0xB6, 0xEC)  # kicker dentro del callout morado oscuro
ON_DARK_CALLOUT_BODY = RGBColor(0xDE, 0xD6, 0xF5)    # texto secundario en fase activa / callout oscuro

FONT = "Inter"

# Constantes de layout — SIEMPRE floats en pulgadas.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_L = 0.6
CONTENT_W = 11.9   # ancho de kicker / headline / body text (margen simétrico ~0.83)
FULL_W = 12.13     # ancho de tablas y fondos internos full-bleed (margen simétrico 0.6)
FOOTER_TOP = 7.12


# ---------------------------------------------------------------------------
# Fundaciones: documento, slide, fondo, texto
# ---------------------------------------------------------------------------

def new_deck():
    """Crea una Presentation en 13.333x7.5in (16:9), sin layouts default."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def add_slide(prs):
    """Agrega un slide en blanco (layout 6) — todo se arma con shapes propios."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG_LIGHT):
    """Rectángulo full-bleed de fondo. BG_LIGHT para contenido, DARK para
    título y cierre. Es shape [0] en todo slide de mt10 — nunca se deja el
    fondo blanco puro de PowerPoint por defecto."""
    return _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=color)


def _run(paragraph, text, size, color=INK, bold=False, italic=False, font=FONT):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


# Alias público — para un párrafo suelto que no encaja en ningún patrón de
# arriba (ej. una quote grande de una sola oración), un build-pptx.py puede
# llamar a `k.run(paragraph, texto, size, ...)` en vez de armar el run a
# mano con la API cruda de python-pptx. Para texto multi-línea o con
# viñetas, usar `bullets_dash` o uno de los patrones documentados.
run = _run


def textbox(slide, left, top, width, height, word_wrap=True, anchor=None):
    """left/top/width/height: floats en pulgadas."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    if anchor is not None:
        tf.vertical_anchor = anchor
    return box, tf


def _shape(slide, shape_type, left, top, width, height, fill=None, line=False):
    """left/top/width/height: floats en pulgadas."""
    shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
# Kicker + headline — cabecera estándar de cualquier slide de contenido.
# Kicker top=0.5in h=0.4in (13pt bold BRAND, MAYÚSCULA), headline top=0.85in
# h=1.15in (26pt bold INK). Sin línea divisoria — mt10 no la usa.
# Devuelve (float) el `top` donde debería empezar el contenido del cuerpo.
# ---------------------------------------------------------------------------

def kicker_headline(slide, kicker, headline, kicker_color=BRAND, content_top=2.30,
                     headline_size=26):
    _, tf = textbox(slide, MARGIN_L, 0.5, CONTENT_W, 0.4)
    _run(tf.paragraphs[0], kicker.upper(), 13, color=kicker_color, bold=True)

    _, tf2 = textbox(slide, MARGIN_L, 0.85, CONTENT_W, 1.15)
    lines = headline.split("\n") if "\n" in headline else [headline]
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        _run(p, line, headline_size, color=INK, bold=True)
    return content_top


# ---------------------------------------------------------------------------
# Footer de marca — bottom-left "MARCA — Sección", bottom-right N de página.
# Siempre en top=7.12in, 9pt MUTED (o ON_DARK_FAINT sobre fondo oscuro).
# ---------------------------------------------------------------------------

def add_speaker_notes(slide, text):
    """Setea el texto de las notas del orador de una slide (panel de notas
    de PowerPoint/Keynote — nunca visible en la slide en sí, ni al
    proyectar en modo presentación estándar).

    Usar para contenido de proceso interno del equipo (gaps pendientes,
    preguntas reservadas para Q&A) que no debe competir visualmente con el
    contenido dirigido a la audiencia, en vez de un `callout_box` visible.
    Ver `pptx-pitch-deck.md`, sección "Notas del orador", para el criterio
    de qué va en notas vs. qué queda en la slide como `callout_box`."""
    slide.notes_slide.notes_text_frame.text = text


def footer_marca(slide, marca, seccion, page_num, on_dark=False):
    color = ON_DARK_FAINT if on_dark else MUTED
    _, tf = textbox(slide, MARGIN_L, FOOTER_TOP, 9.5, 0.3)
    _run(tf.paragraphs[0], f"{marca}  ·  {seccion}", 9, color=color)

    _, tf2 = textbox(slide, 12.30, FOOTER_TOP, 0.6, 0.3)
    _run(tf2.paragraphs[0], str(page_num), 9, color=color)


# ---------------------------------------------------------------------------
# Bullets con guion largo — patrón de contenido de mt10: "—  " en BRAND
# bold seguido del texto en INK. Nivel 1 usa "-  " y tamaño -2pt.
# items: lista de str, o (str, nivel) para sub-items.
# ---------------------------------------------------------------------------

def bullets_dash(tf, items, size=15, color=INK, dash_color=BRAND, space_after=10,
                  bold_level0=False):
    first = True
    for item in items:
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        first = False
        p.level = level
        marker = "—  " if level == 0 else "-  "
        _run(p, marker, size - (2 if level else 0), color=dash_color, bold=True)
        _run(p, text, size - (2 if level else 0), color=color,
             bold=bold_level0 and level == 0)
        p.space_after = Pt(space_after)


def bullets_height(items, width=CONTENT_W, size=15, space_after=10):
    """Estima en pulgadas la altura que va a ocupar un bloque `bullets_dash`
    con estos `items` en un textbox de `width` pulgadas de ancho.

    No hay forma de medir el layout real de PowerPoint sin abrirlo (python-
    pptx no calcula wrapping), así que esto es una aproximación por ancho de
    caracter promedio — deliberadamente conservadora (sobreestima antes que
    subestimar) para dejar margen de sobra si la fuente real (Inter/Inter
    Tight) resulta un poco más ancha que el estimado, o si el archivo se
    abre en una máquina donde la fuente no está instalada y cae a un
    fallback más ancho (ver pptx-pitch-deck.md, sección de fuentes).

    Usar el resultado para posicionar dinámicamente cualquier elemento que
    vaya *debajo* de un bloque de bullets (ej. `roadmap_arrows`) en vez de
    hardcodear un offset fijo que solo funciona para la cantidad de texto
    actual — si el texto crece, el offset fijo se rompe en silencio.
    """
    # Ancho promedio de caracter para Inter/Inter Tight en este tamaño,
    # pulgadas — factor conservador (ancho real promedio ronda 0.46-0.50).
    avg_char_w_in = size * 0.56 / 72.0
    line_h_in = size * 1.28 / 72.0
    marker_chars = 3  # "—  " o "-  "
    total = 0.0
    for item in items:
        text, level = item if isinstance(item, tuple) else (item, 0)
        indent_in = 0.35 if level else 0.0
        avail_w_in = max(1.0, width - indent_in)
        chars_per_line = max(10, int(avail_w_in / avg_char_w_in))
        n_chars = len(text) + marker_chars
        n_lines = max(1, -(-n_chars // chars_per_line))  # ceil sin importar math
        total += n_lines * line_h_in + (space_after / 72.0)
    return total


# ---------------------------------------------------------------------------
# Stat callout — la cifra grande "suelta" (sin card), como "35—45%" o
# "<800 ms" en mt10. Número gigante (30-80pt bold) + label chico debajo
# (13-17pt MUTED). Es el patrón para EL número que domina un slide (o una
# de varias cifras en fila, bajando `big_size`/`width`).
# ---------------------------------------------------------------------------

def stat_callout(slide, big_text, label_lines, left=MARGIN_L, top=2.35,
                  width=5.6, big_size=52, big_color=BRAND, label_size=15.5):
    _, tf = textbox(slide, left, top, width, 1.30)
    _run(tf.paragraphs[0], big_text, big_size, color=big_color, bold=True)

    lines = label_lines if isinstance(label_lines, (list, tuple)) else [label_lines]
    _, tf2 = textbox(slide, left + 0.05, top + 1.15, width - 0.1, 0.9)
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        _run(p, line, label_size, color=MUTED)


# ---------------------------------------------------------------------------
# Stat card — la variante "en caja": rounded-rect blanca con cifra grande +
# label + lista de bullets_dash adentro (como el "12.000.000" de mt10
# slide 2). Usarla cuando el stat necesita contexto/soporte al lado del
# stat_callout suelto, no en vez de él.
# ---------------------------------------------------------------------------

def stat_card(slide, big_text, label, bullet_items=None, left=7.40, top=2.50,
              width=5.35, height=3.55, big_size=28, label_size=13):
    """`label` acepta un str (una línea, comportamiento original) o una
    lista de str (varias líneas, igual que `stat_callout`) — la altura del
    bloque de label y la posición de `bullet_items` se ajustan solas según
    cuántas líneas tenga, así que agregar una segunda línea nunca empuja el
    bloque de bullets a superponerse con el label."""
    card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE)
    card.adjustments[0] = 0.06

    inset_l = left + 0.35
    inset_w = width - 0.65
    # word_wrap=False: a diferencia de stat_callout (que usa el ancho completo
    # del slide), stat_card tiene ~0.7in menos de ancho útil por el padding
    # interno de la card — un big_text que entraría cómodo en un stat_callout
    # puede wrappear acá y superponerse visualmente con el label de abajo. Sin
    # wrap, en el peor caso el número se recorta contra el borde de la card en
    # vez de invadir la línea siguiente (bug real visto en la fila de 4
    # stat_card de "Por qué ahora" — ver tp-ux-ui-designer memory).
    _, tf = textbox(slide, inset_l, top + 0.30, inset_w, 0.50, word_wrap=False)
    _run(tf.paragraphs[0], big_text, big_size, color=BRAND, bold=True)

    label_lines = label if isinstance(label, (list, tuple)) else [label]
    label_line_h = label_size * 0.0205  # ~0.27in a 13pt — estimación conservadora
    label_block_h = label_line_h * len(label_lines) + 0.15
    _, tf2 = textbox(slide, inset_l, top + 0.88, inset_w, label_block_h)
    for i, line in enumerate(label_lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        _run(p, line, label_size, color=MUTED)

    if bullet_items:
        bullets_top = top + 0.88 + label_block_h + 0.20
        bullets_h = max(0.3, height - (bullets_top - top) - 0.15)
        _, tf3 = textbox(slide, inset_l, bullets_top, inset_w, bullets_h)
        bullets_dash(tf3, bullet_items, size=13, space_after=8)
    return card


# ---------------------------------------------------------------------------
# Chip / badge — rounded-rect chico usado suelto (categoría, competidor) o
# en fila (chip_strip, step-strip de sub-sección). Activo = fill BRAND +
# texto blanco bold; inactivo = fill BORDER + texto MUTED.
# ---------------------------------------------------------------------------

def chip(slide, text, left, top, width=2.28, height=0.42, active=False, size=11):
    fill = BRAND if active else BORDER
    text_color = WHITE if active else MUTED
    box = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill)
    box.adjustments[0] = 0.5
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    _run(tf.paragraphs[0], text, size, color=text_color, bold=active)
    return box


def chip_strip(slide, labels, active_index, top=0.5, left=MARGIN_L, chip_width=2.28,
               chip_height=0.42, gap=0.12, size=11):
    """Fila de chips — patrón "paso N de M" de mt10 (ver slides de
    'LA PROPUESTA'). El kicker_headline de ese slide debe indicar
    "SECCIÓN · N DE M" a juego con el chip activo."""
    x = left
    for i, label in enumerate(labels):
        chip(slide, label, x, top, width=chip_width, height=chip_height,
             active=(i == active_index), size=size)
        x = x + chip_width + gap


def chip_card(slide, label, description, left, top, width=5.75, height=1.6,
              active_chip=True, chip_size=11, desc_size=13):
    """Card blanca con un chip/badge (categoría, competidor) arriba y una
    descripción debajo. Composición de `stat_card`-style card + `chip()` —
    patrón para comparar N ítems con una etiqueta corta cada uno (ej. fila
    de competidores), cuando una tabla sería demasiado angosta para el
    texto de cada celda."""
    card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE)
    card.adjustments[0] = 0.06
    chip(slide, label, left + 0.25, top + 0.25, width=min(width - 0.5, 3.2), height=0.38,
         active=active_chip, size=chip_size)
    _, tf = textbox(slide, left + 0.25, top + 0.78, width - 0.50, height - 0.98)
    _run(tf.paragraphs[0], description, desc_size, color=INK)
    return card


# ---------------------------------------------------------------------------
# Callout box — rounded-rect pale con texto itálico. kind="note" (BRAND_PALE,
# neutral/metodológico) o kind="risk" (ORANGE_PALE, riesgo/gap abierto).
# Nunca amarillo — mt10 no usa amarillo para esto, solo morado pálido y
# naranja pálido.
# ---------------------------------------------------------------------------

def callout_box(slide, text, top, kind="note", left=MARGIN_L, width=11.9, height=0.9,
                 size=12.5):
    pale = BRAND_PALE if kind == "note" else ORANGE_PALE
    box = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=pale)
    box.adjustments[0] = 0.10
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _run(tf.paragraphs[0], text, size, color=INK, italic=True)
    return top + height + 0.2


# ---------------------------------------------------------------------------
# Tabla comparativa por shapes — NO usa slide.shapes.add_table: mt10 arma
# las tablas con rectángulos + textboxes para controlar tipografía/spacing
# 1:1 con el resto del deck. Header oscuro (INK) + línea de acento, filas
# alternando blanco / BG_LIGHT, hairline BORDER entre filas.
# rows: lista de listas de str (mismo largo que headers).
# col_widths: lista de floats (inches) que sume ~= width.
# highlight_row: índice de fila a destacar con fill BRAND_PALE (opcional).
# ---------------------------------------------------------------------------

def comparison_table(slide, headers, rows, col_widths, top, left=MARGIN_L, width=FULL_W,
                      row_height=0.60, header_height=0.46, header_size=13,
                      cell_size=13, highlight_row=None):
    # header
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, header_height, fill=INK)
    # cuadra la esquina inferior (el rounded-rect deja las de abajo redondeadas
    # también; un rectángulo fino encima las "cuadra" contra la primera fila)
    _shape(slide, MSO_SHAPE.RECTANGLE, left, top + header_height - 0.04, width, 0.04, fill=INK)

    x = left
    for h, w in zip(headers, col_widths):
        _, tf = textbox(slide, x + 0.14, top, w - 0.14, header_height, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], h, header_size, color=WHITE, bold=True)
        x += w

    y = top + header_height
    for ri, row in enumerate(rows):
        shaded = highlight_row == ri or (highlight_row is None and ri % 2 == 1)
        if shaded:
            _shape(slide, MSO_SHAPE.RECTANGLE, left, y, width, row_height,
                   fill=BG_LIGHT if highlight_row is None else BRAND_PALE)
        # hairline divisor
        _shape(slide, MSO_SHAPE.RECTANGLE, left, y, width, 0.01, fill=BORDER)
        x = left
        for ci, (val, w) in enumerate(zip(row, col_widths)):
            _, tf = textbox(slide, x + 0.14, y, w - 0.14, row_height, anchor=MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            p.word_wrap = True
            _run(p, val, cell_size, color=INK, bold=(ci == 1 and highlight_row is not None))
            x += w
        y += row_height
    return y  # top disponible después de la tabla


# ---------------------------------------------------------------------------
# Step cards — 3 (o N) tarjetas blancas con círculo numerado + título +
# descripción (patrón "Fase 0" de mt10: PoC / Piloto / Benchmark). Usar
# para experimentos, hitos o pasos de un mismo bloque sin orden temporal
# estricto de flujo (para flujo secuencial con flechas, ver roadmap_arrows).
# ---------------------------------------------------------------------------

def step_cards(slide, steps, top=2.30, left=MARGIN_L, card_width=3.90, card_height=3.75,
               gap=0.20):
    x = left
    for i, (title, desc) in enumerate(steps, start=1):
        card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_width, card_height,
                      fill=WHITE)
        card.adjustments[0] = 0.06

        badge = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.30, top + 0.30, 0.55, 0.55,
                       fill=BRAND)
        badge.adjustments[0] = 0.5
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(btf.paragraphs[0], str(i), 20, color=WHITE, bold=True)

        _, tf = textbox(slide, x + 0.30, top + 1.05, card_width - 0.60, 0.85)
        _run(tf.paragraphs[0], title, 16, color=INK, bold=True)

        _, tf2 = textbox(slide, x + 0.30, top + 1.90, card_width - 0.60, card_height - 2.05)
        _run(tf2.paragraphs[0], desc, 12.5, color=MUTED)

        x += card_width + gap
    return top + card_height


# ---------------------------------------------------------------------------
# Roadmap arrows — N tarjetas de fase conectadas por flechas (patrón "Tres
# fases con gates de decisión" de mt10). La fase `active_index` se pinta en
# BRAND (fase actual/próxima); el resto queda blanca (futuras/pasadas).
# phases: lista de (kicker, titulo, subs) donde subs es str, lista de str, o None.
# ---------------------------------------------------------------------------

def roadmap_arrows(slide, phases, active_index=0, top=2.30, left=MARGIN_L,
                    card_width=2.55, card_height=1.75, arrow_width=0.42,
                    arrow_height=0.40, gap_after_card=0.06, gap_after_arrow=0.07):
    x = left
    n = len(phases)
    for i, (kicker, title, subs) in enumerate(phases):
        active = i == active_index
        fill = BRAND if active else WHITE
        kicker_color = ON_DARK_CALLOUT_BODY if active else MUTED
        title_color = WHITE if active else INK
        sub_color = ON_DARK_CALLOUT_BODY if active else MUTED

        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_width, card_height,
               fill=fill).adjustments[0] = 0.08

        _, tf = textbox(slide, x + 0.22, top + 0.18, card_width - 0.40, 0.35)
        _run(tf.paragraphs[0], kicker, 12, color=kicker_color, bold=True)

        _, tf2 = textbox(slide, x + 0.22, top + 0.52, card_width - 0.40, 0.50)
        _run(tf2.paragraphs[0], title, 18, color=title_color, bold=True)

        if subs:
            _, tf3 = textbox(slide, x + 0.22, top + 1.10, card_width - 0.40, card_height - 1.15)
            for j, line in enumerate(subs if isinstance(subs, (list, tuple)) else [subs]):
                p = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
                _run(p, line, 12, color=sub_color)

        x += card_width
        if i < n - 1:
            arrow_x = x + gap_after_card
            arrow_y = top + (card_height - arrow_height) / 2
            _shape(slide, MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, arrow_width, arrow_height,
                   fill=BRAND)
            x = arrow_x + arrow_width + gap_after_arrow
    return top + card_height


# ---------------------------------------------------------------------------
# Slide de título — fondo DARK, barrita de acento (0.5x0.06in BRAND), título
# blanco gigante, subtítulo ON_DARK_SUBTLE, meta-línea ON_DARK_TINT,
# créditos ON_DARK_FAINT.
# ---------------------------------------------------------------------------

def slide_titulo(prs, titulo, subtitulo, meta, creditos):
    s = add_slide(prs)
    set_bg(s, DARK)

    _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.60, 2.55, 0.50, 0.06, fill=BRAND).adjustments[0] = 0.5

    _, tf = textbox(s, MARGIN_L, 2.75, 11.5, 1.30)
    _run(tf.paragraphs[0], titulo, 42, color=WHITE, bold=True)

    _, tf2 = textbox(s, MARGIN_L, 3.75, 11.5, 0.70)
    _run(tf2.paragraphs[0], subtitulo, 18, color=ON_DARK_SUBTLE)

    _, tf3 = textbox(s, MARGIN_L, 5.40, 11.5, 0.40)
    _run(tf3.paragraphs[0], meta, 13, color=ON_DARK_TINT)

    _, tf4 = textbox(s, MARGIN_L, 6.90, 7.0, 0.40)
    _run(tf4.paragraphs[0], creditos, 10.5, color=ON_DARK_FAINT)
    return s


# ---------------------------------------------------------------------------
# Slide de cierre — fondo DARK, quote grande blanca, subtexto opcional
# ON_DARK_SUBTLE, callout kicker+texto en DARK_PURPLE.
# ---------------------------------------------------------------------------

def slide_cierre(prs, quote_lines, ask_kicker, ask_text, subtext_lines=None):
    s = add_slide(prs)
    set_bg(s, DARK)

    _, tf = textbox(s, MARGIN_L, 1.70, 11.80, 1.60)
    lines = quote_lines if isinstance(quote_lines, (list, tuple)) else [quote_lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, line, 32, color=WHITE, bold=True)

    if subtext_lines:
        _, tf2 = textbox(s, MARGIN_L, 3.35, 11.0, 1.30)
        sublines = subtext_lines if isinstance(subtext_lines, (list, tuple)) else [subtext_lines]
        for i, line in enumerate(sublines):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            _run(p, line, 16, color=ON_DARK_SUBTLE)

    _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, 5.00, 11.80, 1.55,
           fill=DARK_PURPLE).adjustments[0] = 0.08
    _, tf3 = textbox(s, 0.95, 5.22, 11.10, 0.40)
    _run(tf3.paragraphs[0], ask_kicker.upper(), 13, color=ON_DARK_CALLOUT_KICKER, bold=True)
    _, tf4 = textbox(s, 0.95, 5.62, 11.10, 0.85)
    _run(tf4.paragraphs[0], ask_text, 14, color=WHITE)
    return s


# ---------------------------------------------------------------------------
# Guardado del deck — SIEMPRE usar `save_deck(prs, path)` en vez de
# `prs.save(path)` directo. Motivo: bug real de python-pptx (verificado en
# 1.0.2, ver `_fix_notes_master_id_lst`) — cualquier deck que haya usado
# `add_speaker_notes()` en alguna slide queda con un .pptx que PowerPoint
# abre bien pero que Keynote rechaza por completo ("can't be imported. The
# file format is invalid.") si se guarda con `prs.save()` a secas.
# ---------------------------------------------------------------------------

def _fix_notes_master_id_lst(prs):
    """Bug real de python-pptx (verificado en 1.0.2, puede seguir presente
    en otras versiones — revisar si esto sigue haciendo falta al upgradear
    la dependencia): acceder a `slide.notes_slide` por primera vez crea el
    part `ppt/notesMasters/notesMaster1.xml` y la relationship en
    `presentation.xml.rels`, pero nunca agrega el elemento
    `<p:notesMasterIdLst>` a `presentation.xml`. PowerPoint tolera esa
    omisión (infiere el notesMaster por relationship type); el importador
    de Keynote la valida estrictamente y rechaza el paquete .pptx ENTERO
    como formato inválido si falta — no solo la slide con notas, el
    archivo completo no abre. Confirmado empíricamente: un .pptx mínimo
    con una nota vía `slide.notes_slide.notes_text_frame.text = ...` fue
    rechazado por Keynote con ese error exacto; el mismo archivo sin notas
    abrió sin problema. Llamar una vez, justo antes de guardar, si el
    deck usó `add_speaker_notes()` en algún lado — no hace nada (no-op) si
    el deck no tiene notas, o si `<p:notesMasterIdLst>` ya está presente."""
    presentation_elm = prs.part._element
    if presentation_elm.find(qn("p:notesMasterIdLst")) is not None:
        return  # ya presente (o python-pptx lo arregló en una version futura)
    notes_master_rId = next(
        (rId for rId, rel in prs.part.rels.items() if rel.reltype == RT.NOTES_MASTER),
        None,
    )
    if notes_master_rId is None:
        return  # el deck no usa notas, nada que arreglar
    lst = etree.SubElement(presentation_elm, qn("p:notesMasterIdLst"))
    nid = etree.SubElement(lst, qn("p:notesMasterId"))
    nid.set(qn("r:id"), notes_master_rId)
    presentation_elm.find(qn("p:sldMasterIdLst")).addnext(lst)


def save_deck(prs, path):
    """Guardar el deck final — SIEMPRE usar esto en vez de `prs.save()`
    directo si el `build-pptx.py` llamó `add_speaker_notes()` en algún
    lado, o Keynote va a rechazar el archivo entero al abrirlo. Ver
    `_fix_notes_master_id_lst` para el motivo. No tiene costo llamarlo
    siempre, incluso en decks sin notas (es no-op en ese caso) — usarlo
    como el único punto de guardado de cualquier `build-pptx.py` nuevo."""
    _fix_notes_master_id_lst(prs)
    prs.save(path)
